import os
import tempfile
import subprocess
import asyncio
from typing import Dict, Any, Optional
from loguru import logger
import PyPDF2
import pdfplumber
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import pytesseract
from PIL import Image
import cv2
import numpy as np
from pathlib import Path

class DocumentProcessor:
    """
    Classe principale pour le traitement et l'extraction de contenu des documents
    """
    
    def __init__(self):
        self.temp_dir = os.environ.get('TEMP_DIR', '/tmp')
        self.supported_formats = {
            'pdf': self._process_pdf,
            'docx': self._process_docx,
            'doc': self._process_doc_legacy,
            'txt': self._process_text,
            'rtf': self._process_rtf,
            'pptx': self._process_pptx,
            'ppt': self._process_ppt_legacy,
            'xlsx': self._process_xlsx,
            'xls': self._process_xls_legacy,
            'csv': self._process_csv,
            'png': self._process_image,
            'jpg': self._process_image,
            'jpeg': self._process_image,
            'tiff': self._process_image,
            'bmp': self._process_image,
            'gif': self._process_image
        }
    
    async def process_file(
        self,
        file_path: str,
        file_type: str,
        extract_text: bool = True,
        perform_ocr: bool = True,
        language: str = 'fr'
    ) -> Dict[str, Any]:
        """
        Traite un fichier et extrait son contenu
        """
        try:
            # Vérification de l'existence du fichier
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"Fichier non trouvé: {file_path}")
            
            # Sélection du processeur approprié
            processor = self.supported_formats.get(file_type.lower())
            if not processor:
                raise ValueError(f"Type de fichier non supporté: {file_type}")
            
            logger.info(f"Traitement du fichier {file_path} avec le processeur {file_type}")
            
            # Exécution du traitement
            result = await processor(
                file_path,
                extract_text=extract_text,
                perform_ocr=perform_ocr,
                language=language
            )
            
            # Ajout des métadonnées de base
            result['metadata'].update({
                'file_size': os.path.getsize(file_path),
                'file_type': file_type,
                'original_filename': os.path.basename(file_path)
            })
            
            return result
            
        except Exception as e:
            logger.error(f"Erreur lors du traitement de {file_path}: {str(e)}")
            raise
    
    async def _process_pdf(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """
        Traite les fichiers PDF
        """
        text_content = ""
        metadata = {}
        
        try:
            # Tentative d'extraction de texte direct
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                metadata['pages'] = len(pdf_reader.pages)
                
                # Extraction des métadonnées PDF
                if pdf_reader.metadata:
                    metadata.update({
                        'title': pdf_reader.metadata.get('/Title', ''),
                        'author': pdf_reader.metadata.get('/Author', ''),
                        'subject': pdf_reader.metadata.get('/Subject', ''),
                        'creator': pdf_reader.metadata.get('/Creator', ''),
                        'creation_date': str(pdf_reader.metadata.get('/CreationDate', ''))
                    })
                
                # Extraction de texte page par page
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_content += f"\n--- Page {page_num + 1} ---\n{page_text}"
                    except Exception as e:
                        logger.warning(f"Erreur extraction page {page_num + 1}: {e}")
            
            # Si peu de texte extrait et OCR demandé, utiliser pdfplumber + OCR
            if (len(text_content.strip()) < 100 and kwargs.get('perform_ocr', True)):
                logger.info("Peu de texte extrait, tentative avec pdfplumber et OCR")
                text_content += await self._pdf_ocr_fallback(file_path, kwargs.get('language', 'fr'))
            
        except Exception as e:
            logger.error(f"Erreur traitement PDF: {e}")
            # Fallback vers OCR si extraction directe échoue
            if kwargs.get('perform_ocr', True):
                text_content = await self._pdf_ocr_fallback(file_path, kwargs.get('language', 'fr'))
        
        return {
            'text_content': text_content.strip(),
            'metadata': metadata
        }
    
    async def _pdf_ocr_fallback(self, file_path: str, language: str = 'fr') -> str:
        """
        OCR de secours pour les PDF
        """
        try:
            # Conversion PDF vers images puis OCR
            output_dir = tempfile.mkdtemp()
            cmd = [
                'pdftoppm',
                file_path,
                f'{output_dir}/page',
                '-png'
            ]
            
            process = await asyncio.create_subprocess_exec(
                *cmd,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            await process.communicate()
            
            # OCR sur chaque image générée
            ocr_text = ""
            for img_file in sorted(Path(output_dir).glob('*.png')):
                page_text = await self._ocr_image(str(img_file), language)
                if page_text.strip():
                    ocr_text += f"\n--- OCR Page ---\n{page_text}"
            
            # Nettoyage
            import shutil
            shutil.rmtree(output_dir)
            
            return ocr_text
            
        except Exception as e:
            logger.error(f"Erreur OCR PDF: {e}")
            return ""
    
    async def _process_docx(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """
        Traite les fichiers DOCX
        """
        try:
            doc = Document(file_path)
            
            # Extraction du texte
            text_content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            # Extraction des métadonnées
            metadata = {
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'images': len(doc.inline_shapes)
            }
            
            # Métadonnées du document
            if hasattr(doc.core_properties, 'title'):
                metadata.update({
                    'title': doc.core_properties.title or '',
                    'author': doc.core_properties.author or '',
                    'subject': doc.core_properties.subject or '',
                    'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                    'modified': str(doc.core_properties.modified) if doc.core_properties.modified else ''
                })
            
            return {
                'text_content': text_content,
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Erreur traitement DOCX: {e}")
            raise
    
    async def _process_doc_legacy(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """
        Traite les fichiers DOC legacy via LibreOffice
        """
        return await self._process_with_libreoffice(file_path, 'doc')
    
    async def _process_with_libreoffice(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """
        Convertit et traite les fichiers via LibreOffice
        """
        try:
            output_dir = tempfile.mkdtemp()
            
            # Conversion vers format moderne via LibreOffice
            cmd = [
                'libreoffice',
                '--headless',
                '--convert-to',
                'txt',
                '--outdir',
                output_dir,
                file_path
            ]
            
            process = await asyncio.create_subprocess_exec(
                *cmd,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            
            stdout, stderr = await process.communicate()
            
            # Lecture du fichier converti
            txt_file = Path(output_dir) / (Path(file_path).stem + '.txt')
            text_content = ""
            
            if txt_file.exists():
                with open(txt_file, 'r', encoding='utf-8') as f:
                    text_content = f.read()
            
            # Nettoyage
            import shutil
            shutil.rmtree(output_dir)
            
            return {
                'text_content': text_content,
                'metadata': {'converted_via': 'libreoffice'}
            }
            
        except Exception as e:
            logger.error(f"Erreur LibreOffice: {e}")
            raise
    
    async def _process_image(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """
        Traite les images avec OCR
        """
        if not kwargs.get('perform_ocr', True):
            return {'text_content': '', 'metadata': {'type': 'image', 'ocr_skipped': True}}
        
        return await self._ocr_image_detailed(file_path, kwargs.get('language', 'fr'))
    
    async def _ocr_image_detailed(self, file_path: str, language: str = 'fr') -> Dict[str, Any]:
        """
        OCR détaillé d'une image avec pré-traitement
        """
        try:
            # Chargement de l'image
            img = cv2.imread(file_path)
            if img is None:
                raise ValueError("Impossible de charger l'image")
            
            # Pré-traitement pour améliorer l'OCR
            img = await self._preprocess_image_for_ocr(img)
            
            # OCR avec Tesseract
            lang_map = {'fr': 'fra', 'en': 'eng'}
            tesseract_lang = lang_map.get(language, 'fra')
            
            # Configuration Tesseract
            config = '--oem 3 --psm 6'
            
            text_content = pytesseract.image_to_string(
                img,
                lang=tesseract_lang,
                config=config
            )
            
            # Extraction des données de confiance
            data = pytesseract.image_to_data(
                img,
                lang=tesseract_lang,
                config=config,
                output_type=pytesseract.Output.DICT
            )
            
            # Calcul de la confiance moyenne
            confidences = [int(conf) for conf in data['conf'] if int(conf) > 0]
            avg_confidence = sum(confidences) / len(confidences) if confidences else 0
            
            metadata = {
                'ocr_confidence': avg_confidence,
                'image_dimensions': img.shape[:2],
                'detected_text_blocks': len([conf for conf in confidences if conf > 50])
            }
            
            return {
                'text_content': text_content,
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Erreur OCR image: {e}")
            return {'text_content': '', 'metadata': {'ocr_error': str(e)}}
    
    async def _preprocess_image_for_ocr(self, img):
        """
        Pré-traitement d'image pour améliorer l'OCR
        """
        # Conversion en niveaux de gris
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        
        # Réduction du bruit
        denoised = cv2.medianBlur(gray, 3)
        
        # Amélioration du contraste
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8,8))
        enhanced = clahe.apply(denoised)
        
        # Binarisation adaptative
        binary = cv2.adaptiveThreshold(
            enhanced, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
        )
        
        return binary
    
    async def _ocr_image(self, file_path: str, language: str = 'fr') -> str:
        """
        OCR simple d'une image
        """
        try:
            lang_map = {'fr': 'fra', 'en': 'eng'}
            tesseract_lang = lang_map.get(language, 'fra')
            
            return pytesseract.image_to_string(
                Image.open(file_path),
                lang=tesseract_lang
            )
        except Exception as e:
            logger.error(f"Erreur OCR simple: {e}")
            return ""
    
    # Méthodes pour autres formats
    async def _process_text(self, file_path: str, **kwargs) -> Dict[str, Any]:
        with open(file_path, 'r', encoding='utf-8') as f:
            return {'text_content': f.read(), 'metadata': {}}
    
    async def _process_pptx(self, file_path: str, **kwargs) -> Dict[str, Any]:
        prs = Presentation(file_path)
        text_content = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content += shape.text + "\n"
        
        return {
            'text_content': text_content,
            'metadata': {'slides': len(prs.slides)}
        }
    
    async def _process_xlsx(self, file_path: str, **kwargs) -> Dict[str, Any]:
        wb = load_workbook(file_path)
        text_content = ""
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_content += f"\n--- Feuille: {sheet_name} ---\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text_content += row_text + "\n"
        
        return {
            'text_content': text_content,
            'metadata': {'sheets': len(wb.sheetnames)}
        }
    
    # Méthodes pour formats legacy
    async def _process_rtf(self, file_path: str, **kwargs) -> Dict[str, Any]:
        return await self._process_with_libreoffice(file_path, 'rtf')
    
    async def _process_ppt_legacy(self, file_path: str, **kwargs) -> Dict[str, Any]:
        return await self._process_with_libreoffice(file_path, 'ppt')
    
    async def _process_xls_legacy(self, file_path: str, **kwargs) -> Dict[str, Any]:
        return await self._process_with_libreoffice(file_path, 'xls')
    
    async def _process_csv(self, file_path: str, **kwargs) -> Dict[str, Any]:
        import csv
        text_content = ""
        
        with open(file_path, 'r', encoding='utf-8') as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                text_content += "\t".join(row) + "\n"
        
        return {
            'text_content': text_content,
            'metadata': {'format': 'csv'}
        }